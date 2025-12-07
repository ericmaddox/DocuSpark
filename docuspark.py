import os
from pathlib import Path
from PIL import Image
import pytesseract
import pdfplumber
from docx import Document
from pptx import Presentation
import pypandoc
import io

# --------------------------
# Configuration
# --------------------------
RAW_DIR = "data"
OUTPUT_DIR = "clean_md"
SUPPORTED_TYPES = {".pdf", ".docx", ".pptx", ".txt", ".html", ".htm", ".rtf"}

# --------------------------
# Helper functions
# --------------------------
def ensure_dir(path):
    os.makedirs(path, exist_ok=True)

def describe_image(image_path):
    """
    Generate a textual description of an image (OCR-based).
    """
    try:
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img, lang='eng')
        if not text.strip():
            text = "No readable text detected. Consider using a captioning model."
        return text.strip()
    except Exception as e:
        return f"Failed to describe image: {e}"

def add_images_to_markdown(md, image_paths):
    """
    Append images and their descriptions to the Markdown content.
    """
    if not image_paths:
        return md

    lines = [md, "\n\n## Extracted Images\n"]
    for idx, img_path in enumerate(image_paths, start=1):
        rel_img_path = os.path.relpath(img_path, start=OUTPUT_DIR)
        lines.append(f"![Figure {idx}]({rel_img_path})")
        description = describe_image(img_path)
        lines.append(f"**Description:** {description}\n")
    return "\n".join(lines)

# --------------------------
# File converters
# --------------------------
def convert_pdf(input_path):
    """Extract text and images from PDF."""
    text_parts = []
    images = []
    try:
        with pdfplumber.open(input_path) as pdf:
            for page_idx, page in enumerate(pdf.pages, start=1):
                text_parts.append(page.extract_text() or "")
                for img_idx, img in enumerate(page.images, start=1):
                    try:
                        bbox = (img["x0"], img["top"], img["x1"], img["bottom"])
                        im = page.crop(bbox).to_image(resolution=300).original
                        img_name = f"img_{page_idx}_{img_idx}.png"
                        images.append((im, img_name))
                    except:
                        continue
    except Exception as e:
        print(f"[PDF] Failed on {input_path}: {e}")

    return "\n\n".join(text_parts), images

def convert_docx(input_path):
    """Extract text from DOCX and images."""
    text_parts = []
    images = []
    try:
        doc = Document(input_path)
        for para in doc.paragraphs:
            text_parts.append(para.text)
        for idx, rel in enumerate(doc.part.rels.values(), start=1):
            if "image" in rel.target_ref:
                try:
                    img_bytes = rel.target_part.blob
                    im = Image.open(io.BytesIO(img_bytes))
                    img_name = f"img_{idx}.png"
                    images.append((im, img_name))
                except:
                    continue
    except Exception as e:
        print(f"[DOCX] Failed on {input_path}: {e}")
    return "\n".join(text_parts), images

def convert_pptx(input_path):
    """Extract text and images from PPTX."""
    text_parts = []
    images = []
    try:
        prs = Presentation(input_path)
        for slide_idx, slide in enumerate(prs.slides, start=1):
            text_parts.append(f"# Slide {slide_idx}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
                if shape.shape_type == 13:  # picture
                    try:
                        img = shape.image
                        im = Image.open(io.BytesIO(img.blob))
                        img_name = f"img_slide{slide_idx}.png"
                        images.append((im, img_name))
                    except:
                        continue
    except Exception as e:
        print(f"[PPTX] Failed on {input_path}: {e}")
    return "\n\n".join(text_parts), images

def convert_generic(input_path):
    """Fallback for TXT, HTML, RTF using Pandoc."""
    try:
        text = pypandoc.convert_file(input_path, "md")
        return text, []
    except Exception as e:
        print(f"[Pandoc] Failed on {input_path}: {e}")
        return None, []

def convert_file(input_path):
    ext = os.path.splitext(input_path)[1].lower()
    if ext == ".pdf":
        return convert_pdf(input_path)
    elif ext == ".docx":
        return convert_docx(input_path)
    elif ext == ".pptx":
        return convert_pptx(input_path)
    elif ext in {".txt", ".html", ".htm", ".rtf"}:
        return convert_generic(input_path)
    else:
        return None, []

# --------------------------
# Pipeline
# --------------------------
def run_pipeline():
    print("\n🔍 Starting DocuSpark ingestion…")

    if not os.path.exists(RAW_DIR):
        print(f"❌ Folder '{RAW_DIR}' does not exist.")
        return

    for root, dirs, files in os.walk(RAW_DIR):
        for file in sorted(files):
            ext = os.path.splitext(file)[1].lower()
            if ext not in SUPPORTED_TYPES:
                continue

            input_path = os.path.join(root, file)
            relative_path = os.path.relpath(root, RAW_DIR)
            output_dir = os.path.join(OUTPUT_DIR, relative_path)
            ensure_dir(output_dir)

            md_filename = os.path.splitext(file)[0] + ".md"
            md_output_path = os.path.join(output_dir, md_filename)

            print(f"📄 Processing: {input_path}")
            markdown_text, images = convert_file(input_path)

            if markdown_text is None:
                print(f"⚠️ Skipped: {input_path}")
                continue

            # Save images and update markdown
            img_paths = []
            for im, img_name in images:
                img_folder = os.path.join(output_dir, "images")
                ensure_dir(img_folder)
                img_path = os.path.join(img_folder, img_name)
                try:
                    im.save(img_path)
                    img_paths.append(img_path)
                except Exception as e:
                    print(f"⚠️ Failed to save image {img_name}: {e}")

            markdown_text = add_images_to_markdown(markdown_text, img_paths)

            # Write final markdown
            with open(md_output_path, "w", encoding="utf-8") as f:
                f.write(markdown_text)

    print("\n✅ Done! Markdown files with image descriptions created in 'clean_md/'\n")

# --------------------------
# Entry point
# --------------------------
if __name__ == "__main__":
    run_pipeline()
